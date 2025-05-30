from datetime import datetime
import gc
import hashlib
import json
import logging
import pytz
import streamlit as st
import os
import io
import pdfplumber
import openai
import pinecone
import time
import unicodedata
import re
from azure.storage.blob import BlobServiceClient
from docx import Document
from pptx import Presentation
import pandas as pd
import gspread
from oauth2client.service_account import ServiceAccountCredentials
from dotenv import load_dotenv

# Load .env variables
load_dotenv()
logs = []

# # ✅ Load environment variables
# AZURE_CONNECTION_STRING = os.getenv("AZURE_CONNECTION_STRING")
# AZURE_CONTAINER_NAME = os.getenv("AZURE_CONTAINER_NAME")
# PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
# PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
# OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
# GOOGLE_SHEET_CREDENTIALS_PATH = os.getenv("CREDENTIALS_PATH")

# ✅ Load environment variables from Streamlit Secrets
AZURE_CONNECTION_STRING = st.secrets["AZURE_CONNECTION_STRING"]
AZURE_CONTAINER_NAME = st.secrets["AZURE_CONTAINER_NAME"]
PINECONE_API_KEY = st.secrets["PINECONE_API_KEY"]
PINECONE_INDEX_NAME = st.secrets["PINECONE_INDEX_NAME"]
OPENAI_API_KEY = st.secrets["OPENAI_API_KEY"]

# ✅ Initialize Azure Blob & Pinecone
blob_service_client = BlobServiceClient.from_connection_string(AZURE_CONNECTION_STRING)
container_client = blob_service_client.get_container_client(AZURE_CONTAINER_NAME)

pinecone_client = pinecone.Pinecone(api_key=PINECONE_API_KEY)
index = pinecone_client.Index(PINECONE_INDEX_NAME)

openai.api_key = OPENAI_API_KEY

# #✅ Google Sheets setup
# scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
# creds = ServiceAccountCredentials.from_json_keyfile_name(GOOGLE_SHEET_CREDENTIALS_PATH, scope)
#✅ Setup Google Sheets with credentials from Streamlit Secrets
credentials_info = json.loads(st.secrets["google_sheets"]["credentials"])  # Parse the JSON credentials string
creds = ServiceAccountCredentials.from_json_keyfile_dict(credentials_info, 
    ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"])
client = gspread.authorize(creds)
sheet = client.open("VectorizationLogs").sheet1  # Replace with your Google Sheet name

# ✅Senitize path of files
def sanitize_file_path(file_path):
    """Sanitize file path for Pinecone."""
    sanitized = unicodedata.normalize("NFKD", file_path).encode("ascii", "ignore").decode("ascii")
    sanitized = re.sub(r"[^a-zA-Z0-9_/.-]", "_", sanitized)
    return sanitized

# ✅GENERATE VECTOR ID USING HASING 
def generate_md5_vector_id(file_path, page_num, chunk_num):
    """Generate an MD5 hash-based vector ID using only the file path, appending page and chunk info."""
    # Hash only the file path
    md5_hash = hashlib.md5(file_path.encode()).hexdigest()
    
    # Return the vector ID in the format: <hash>-page<page_num>-chunk<chunk_num>
    return f"{md5_hash}-page{page_num}-chunk{chunk_num}"

# ✅Extract text from PPT, DOC, PDF 
def extract_text_from_file(file_path, blob_data):
    """Extract text from PDF, DOCX, or PPTX files."""
    all_pages_data = []

    if file_path.endswith(".pdf"):
        try:
            with pdfplumber.open(io.BytesIO(blob_data)) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text()
                    if page_text:
                        all_pages_data.append({"page_num": page_num, "text": page_text})
                return all_pages_data
        except:
            logs.append([file_path, "N/A", "Failure", "Currepted File"])
            all_pages_data=None
            return all_pages_data
            

    elif file_path.endswith(".docx"):
        try:
            with io.BytesIO(blob_data) as doc_stream:
                doc = Document(doc_stream)
                words_per_page = 250
                word_buffer = []
                page_num = 1

                for para in doc.paragraphs:
                    words = para.text.strip().split()
                    if words:
                        word_buffer.extend(words)

                        if len(word_buffer) >= words_per_page:
                            all_pages_data.append({"page_num": page_num, "text": " ".join(word_buffer)})
                            word_buffer = []
                            page_num += 1

                if word_buffer:
                    all_pages_data.append({"page_num": page_num, "text": " ".join(word_buffer)})
                    return all_pages_data
        except:
            logs.append([file_path, "N/A", "Failure", "Currepted File"])
            all_pages_data=None
            return all_pages_data

    elif file_path.endswith(".pptx"):
        try:
          with io.BytesIO(blob_data) as ppt_stream:
            ppt = Presentation(ppt_stream)
            for page_num, slide in enumerate(ppt.slides, start=1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                if slide_text.strip():
                    all_pages_data.append({"page_num": page_num, "text": slide_text.strip()})
            return all_pages_data
        except Exception as e:
            print(f"⚠️ Error reading .pptx: {file_path} - {e}")
            logs.append([file_path, "N/A", "Failure", "Currepted File"])
            all_pages_data = None
            return all_pages_data    
    

# Chunk the extracted text into check✅
def chunk_text_for_openai(text, max_tokens=4096):
    tokens_per_word = 1.3  # Approximate tokens per word
    words = text.split()

    chunk = []
    chunks = []
    token_count = 0

    for word in words:
        chunk.append(word)
        token_count += tokens_per_word  # 1 word = ~1.3 tokens

        if token_count > max_tokens:
            chunks.append(" ".join(chunk))
            chunk = [word]
            token_count = tokens_per_word

    if chunk:
        chunks.append(" ".join(chunk))

    return chunks

# Log Vectorization in Google sheet
def log_to_google_sheets_batch(logs):
    """Batch logging to Google Sheets with only unique file names and timestamps."""
    if not logs:
        st.warning("⚠️ No logs to write to Google Sheets.")
        return
    
    try:
        # Extract unique file names from logs
        unique_files = list(set(log[0] for log in logs))  # Only file names
        
        # Get the current timestamp for logging
        #current_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        # Define the time zone (e.g., UTC)
        timezone = pytz.timezone('UTC')

        # Get the current time in the specified time zone
        current_time = datetime.now(timezone).strftime("%Y-%m-%d %H:%M:%S")
        
        # Prepare formatted logs with timestamp
        formatted_logs = [[file_name, "Success", current_time] for file_name in unique_files]  # Include timestamp
        
        # Print logs for debugging
        st.write("📝 Logs Ready to Write:", formatted_logs)

        # Write logs to the sheet
        sheet.append_rows(formatted_logs, value_input_option="RAW")

        st.success(f"✅ Logs successfully written for {len(unique_files)} files!")

    except Exception as e:
        error_msg = f"Google Sheets Logging Error: {e}"
        print(error_msg)  # Print error to console
        st.error(error_msg)  # Display error in Streamlit UI
   
def is_file_processed(file_path):
    """Check if the file is already stored in Pinecone using the MD5 hash-based vector ID."""
    sanitized_file_path = sanitize_file_path(file_path)  # Sanitize the file path to avoid issues with special characters

    try:
        # Check for vectors for the first page and chunk (using the MD5 hash of the file path)
        chunk_id = generate_md5_vector_id(sanitized_file_path, page_num=1, chunk_num=0)  # Example: page 1, chunk 0
        
        # Fetch vectors for this chunk ID from Pinecone
        existing_vectors = index.fetch([chunk_id])  # Use the generated vector ID to check for existence
        
        # Return True if vectors are found, meaning the file has been processed
        #return bool(existing_vectors.get("vectors"))  # If vectors are found, return True
        return bool(existing_vectors.vectors)  # Instead of existing_vectors.get("vecto

    except Exception as e:
        print(f"Pinecone Query Error: {e}")
        return False     
 
       
# STORE IN PINECONE✅
def store_vectors(file_path, extracted_data):
    """Store vectorized text in Pinecone with chunking per page."""
    sanitized_file_path = sanitize_file_path(file_path)
    blob_client = container_client.get_blob_client(file_path)

    all_chunks = []
    metadata_list = []
    

    try:
        # ✅ Ensure the file is not already processed by checking its status
        if is_file_processed(file_path):
            logging.info(f"File {file_path} already processed.")
            return

        # ✅ Process extracted data in smaller chunks to avoid memory overload
        for page_data in extracted_data:
            # Chunk page text
            chunks = chunk_text_for_openai(page_data["text"], max_tokens=8000)  # Safe buffer for 8191 token limit
            for i, chunk in enumerate(chunks):
                # Use MD5-based vector ID
                chunk_id = generate_md5_vector_id(sanitized_file_path, page_data["page_num"], i)
                metadata = {
                    "text": chunk,
                    "type": "text",
                    "page_number": page_data["page_num"],
                    "title": sanitized_file_path,
                    "link": blob_client.url
                }
                all_chunks.append(chunk)
                metadata_list.append((chunk_id, metadata))

        # ✅ Embed chunks in batches to avoid timeout
        batch_size = 50  # Process 50 chunks at a time
        for start_idx in range(0, len(all_chunks), batch_size):
            batch_chunks = all_chunks[start_idx:start_idx + batch_size]
            batch_metadata = metadata_list[start_idx:start_idx + batch_size]

            response = openai.Embedding.create(input=batch_chunks, model="text-embedding-3-large")
            batch_vectors = []
            for i, (chunk_id, metadata) in enumerate(batch_metadata):
                batch_vectors.append({
                    "id": chunk_id,
                    "values": response["data"][i]["embedding"],
                    "metadata": metadata
                })
                logs.append([file_path, metadata["page_number"], "Success", ""])

            # ✅ Store batch in Pinecone
            index.upsert(vectors=batch_vectors)

            # Clear memory after each batch to prevent overload
            del batch_vectors
            gc.collect()  # Garbage collect to free up memory

            # Optional: Log status after each batch
            logging.info(f"Processed batch {start_idx // batch_size + 1}/{len(all_chunks) // batch_size + 1}")
            
    except openai.error.APIConnectionError as e:
        logging.error(f"APIConnectionError while processing {file_path}: {e}")
        logs.append([file_path, "N/A", "Failure", str(e)])
        time.sleep(5)  # Retry after 5 seconds

    except openai.error.RateLimitError:
        logging.error(f"Rate limit exceeded for {file_path}. Retrying...")
        logs.append([file_path, "N/A", "Failure", "Rate Limit Exceeded"])
        time.sleep(60)  # Retry after 1 minute

    except Exception as e:
        logging.error(f"Unexpected error processing {file_path}: {e}")
        logs.append([file_path, "N/A", "Failure", str(e)])

    finally:
        # ✅ Log results to Google Sheets or a similar service
        if logs:
            log_to_google_sheets_batch(logs)
        logging.info(f"Finished processing {file_path}")
        
# MAIN function that call all sub function and do vectorization✅   
def process_new_files():
    """Check for new files, extract text, vectorize, and store logs."""
    blob_list = container_client.list_blobs()
    new_files = []

    for blob in blob_list:
        file_path = blob.name.strip()

        if file_path.endswith((".pdf", ".docx", ".pptx")) and not is_file_processed(file_path):
            blob_client = container_client.get_blob_client(file_path)
            blob_data = blob_client.download_blob().readall()

            extracted_data = extract_text_from_file(file_path, blob_data)
            if extracted_data is not None:
                store_vectors(file_path, extracted_data)
    
                new_files.append(file_path)
    return new_files

# ✅ Streamlit UI
st.title("📄 Streamlit Auto-Vectorization")
st.write("Automatically process & vectorize documents from Azure Blob Storage.")

# ✅ Check for new files first
blob_list = list(container_client.list_blobs())  # Convert iterator to list
st.subheader("📂 Files Pending Vectorization")

# Get list of all eligible files from Azure
all_files = [blob.name.strip() for blob in blob_list if blob.name.endswith((".pdf", ".docx", ".pptx"))]

# Filter to only those not yet vectorized
files_to_vectorize = [file for file in all_files if not is_file_processed(file)]

if files_to_vectorize:
    st.write(f"🔍 {len(files_to_vectorize)} new files found:")
    st.write("+++++++++++++++++++++++++++++++++++")
    st.info(files_to_vectorize)
else:
    st.success("✅ All eligible files are already vectorized!")
if files_to_vectorize:
    with st.spinner("🔄 Processing all new files..."):
        processed_files = process_new_files()
        st.success(f"✅ Vectorization completed for {len(processed_files)} new files!")
else:
    st.warning("⚠️ No new file is detected.")


# ✅ Display Google Sheets Log
st.subheader("📊 Vectorization Logs")
log_data = sheet.get_all_values()
df = pd.DataFrame(log_data, columns=["File Name", "Page Number", "Status"])
st.dataframe(df)
